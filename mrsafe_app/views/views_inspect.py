import base64
import os
import re
import traceback
from dotenv import load_dotenv

from django.conf import settings
from django.shortcuts import render, redirect
from django.utils.timezone import now
from django.views.decorators.csrf import csrf_exempt
from django.core.files.storage import FileSystemStorage
from openai import OpenAI
from django.contrib.auth import get_user_model

from ..forms import HazardForm, RecommendationForm
from ..models import SafetyObservation,CoinActivity,UserCoinBalance

client = OpenAI(api_key=settings.OPENAI_API_KEY)


def parse_markdown_sections(markdown_text):
    sections = []
    raw_sections = [s.strip() for s in markdown_text.split("##") if s.strip()]
    bullet_pattern = re.compile(r"^(\s*[-*•]|\s*\d+\.)\s+")

    for section in raw_sections:
        lines = section.splitlines()
        title = lines[0].strip() if lines else "Untitled"
        body_lines = lines[1:] if len(lines) > 1 else []
        bullet_points = [re.sub(bullet_pattern, "", line.strip()) for line in body_lines if bullet_pattern.match(line.strip())]
        sections.append({"title": title, "bullets": bullet_points})
    return sections



@csrf_exempt
def safety_image_test(request):
    context = {}

    if not settings.OPENAI_API_KEY:
        context["error"] = "OpenAI API key is missing"
        return render(request, "mrsafe/inspect/inspect.html", context)

    if request.method == "POST" and request.FILES.get("photo"):
        image_file = request.FILES["photo"]

        try:
            if image_file.size > 5 * 1024 * 1024:
                raise ValueError("Image size exceeds 5MB")

            # Save uploaded photo to media/uploads
            fs = FileSystemStorage(location="media/uploads", base_url="/media/uploads/")
            filename = fs.save(image_file.name, image_file)
            uploaded_url = fs.url(filename)
            context["photo_url"] = uploaded_url

            with fs.open(filename, "rb") as img:
                base64_img = base64.b64encode(img.read()).decode("utf-8")

            # Debug: content-type and input size
            content_type = image_file.content_type or "image/jpeg"
            image_ext = content_type.split("/")[-1]
            print("[DEBUG] Content-Type:", content_type)
            print("[DEBUG] Base64 Length:", len(base64_img))

            prompt = """
**Safety Inspection Analysis Request**


You are an AI safety inspection assistant. Based on the provided input (such as a photo or observation), generate clearly separated and numbered hazard reports. 

Each hazard must be output using the following exact structure. Do not combine or merge hazards. Each block must be self-contained.

For every identified issue, follow **this markdown format**:

---

### Hazard #{{ n }}
- **Title**: [Short title of hazard]
- **Severity**: Low / Medium / High / Critical
- **Description**: [Detailed explanation of the hazard and its impact]
- **OSHA Reference**: [e.g., 1926.501(b)(1)] or "N/A" if not applicable

#### Recommendation
- **Action**: [Corrective or preventive measure to eliminate or reduce the hazard]

---

### ✅ Example Output

```markdown
### Hazard #1
- **Title**: Missing Guardrails on Elevated Platform
- **Severity**: High
- **Description**: The platform at the construction site lacks guardrails on two sides, increasing the risk of falls from height.
- **OSHA Reference**: 1926.501(b)(1)

#### Recommendation
- **Action**: Install compliant guardrails (at least 42 inches high) on all open sides of the platform immediately.

---

### Hazard #2
- **Title**: Blocked Emergency Exit
- **Severity**: Critical
- **Description**: The designated emergency exit is obstructed by equipment, making it inaccessible during an emergency.
- **OSHA Reference**: 1910.37(a)(3)

#### Recommendation
- **Action**: Clear the emergency exit path and place visible signage to ensure constant accessibility.


---

Repeat the structure for each hazard found in the image. Do not list hazards or recommendations in separate sections. Always keep each hazard and its recommendation grouped.

Format your output in clean markdown as shown above.
"""

            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "You are a workplace safety AI assistant."},
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/{image_ext};base64,{base64_img}",
                                    "detail": "high"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=1500,
                temperature=0.3,
            )

            print("[DEBUG] GPT Response Object:", response)

            if not response.choices or not response.choices[0].message.content:
                raise ValueError("GPT response missing expected content.")

            analysis = response.choices[0].message.content
            clean_analysis = analysis.replace("```markdown", "").replace("```", "").strip()
            sections = parse_markdown_sections(clean_analysis)
            context["sections"] = sections

            # Extract hazard and recommendation summary
            hazard_lines = []
            reco_lines = []
            for section in sections:
                title = section["title"].strip()
                for bullet in section.get("bullets", []):
                    if "severity" in bullet.lower():
                        hazard_lines.append(f"{title}: {bullet}")
                    elif any(x in title.lower() for x in ["recommendation", "action"]):
                        reco_lines.append(f"{title}: {bullet}")

            # Save to DB
            user = request.user if request.user.is_authenticated else get_user_model().objects.filter(is_active=True).first()
            if not user:
                context["error"] = "No valid user found to assign observation."
                return render(request, "mrsafe/inspect/inspect.html", context)
            
            # ✅ Get or create user (already in your code)
            user = request.user if request.user.is_authenticated else get_user_model().objects.filter(is_active=True).first()
            if not user:
                context["error"] = "No valid user found to assign observation."
                return render(request, "mrsafe/inspect/inspect.html", context)

            # ✅ Try to fetch the "AI Safety Image Analysis" coin activity
            try:
                ai_analysis_activity = CoinActivity.objects.get(name="ai_image_analysis", is_active=True)
            except CoinActivity.DoesNotExist:
                context["error"] = "❌ AI Image Analysis coin activity is not defined. Please contact admin."
                return render(request, "mrsafe/inspect/inspect.html", context)

            # ✅ Get or create user coin balance
            user_balance, _ = UserCoinBalance.objects.get_or_create(user=user)

            # ✅ Check balance
            if user_balance.balance < (ai_analysis_activity.coin_amount or 0):
                context["error"] = "❌ You don't have enough coins to analyze this photo."
                return render(request, "mrsafe/inspect/inspect.html", context)

            # ✅ Deduct coins and log transaction
            success = user_balance.update_balance(
                ai_analysis_activity.coin_amount,
                "spend",
                ai_analysis_activity
            )

            if not success:
                context["error"] = "❌ Failed to deduct coins. Please try again later."
                return render(request, "mrsafe/inspect/inspect.html", context)

            SafetyObservation.objects.create(
                photo=image_file,
                hazard_description="\n".join(hazard_lines),
                recommendations="\n".join(reco_lines),
                created_by=user
            )

            # Forms for UI
            hazard_forms = []
            recommendation_forms = []

            for section in sections:
                title = section["title"].lower()
                for bullet in section["bullets"]:
                    if "severity" in bullet.lower():
                        severity_match = re.search(r"severity:\s*(\w+)", bullet, re.I)
                        osha_match = re.search(r"osha.*?:\s*(.+)", bullet, re.I)
                        hazard_forms.append(HazardForm(initial={
                            "title": section["title"],
                            "description": bullet,
                            "severity": severity_match.group(1) if severity_match else "Medium",
                            "oshas": osha_match.group(1) if osha_match else "",
                        }))
                    elif any(k in title for k in ["recommendation", "action"]):
                        recommendation_forms.append(RecommendationForm(initial={
                            "title": section["title"],
                            "action": bullet
                        }))

            context["hazard_forms"] = hazard_forms
            context["recommendation_forms"] = recommendation_forms
            context["message"] = "✅ Safety observation successfully saved."

        except ValueError as ve:
            context["error"] = f"Validation Error: {str(ve)}"
        except Exception as e:
            context["error"] = f"Analysis Error: {str(e)}"
            if settings.DEBUG:
                context["error_details"] = traceback.format_exc()

    context["now"] = now()
    return render(request, "mrsafe/inspect/inspect.html", context)



def calculate_safety_score(analysis):
    """Calculate safety score based on hazard severity in analysis"""
    critical = analysis.count("Severity: Critical")
    high = analysis.count("Severity: High")
    medium = analysis.count("Severity: Medium")
    low = analysis.count("Severity: Low")
    
    score = 100 - (critical * 20 + high * 10 + medium * 5 + low * 2)
    return max(30, min(100, score))


# mrsafe_app/views/views_inspect.py



def inspect(request):
    uploaded_image_url = None

    if request.method == "POST" and request.FILES.get("safety_image"):
        image_file = request.FILES["safety_image"]
        fs = FileSystemStorage(location="media/uploads", base_url="/media/uploads/")
        filename = fs.save(image_file.name, image_file)
        uploaded_image_url = fs.url(filename)

    return render(request, "mrsafe/inspect/inspect.html", {
        "uploaded_image_url": uploaded_image_url,
    })
    


from django.contrib.auth.decorators import login_required
from django.core.files.base import ContentFile
from django.shortcuts import redirect
import os
from django.conf import settings
from ..models import SafetyObservation  # adjust import if needed

@login_required
def save_observation(request):
    if request.method == 'POST':
        photo_url = request.POST.get('photo')
        hazard_description = request.POST.get('hazard_description', '')
        recommendations = request.POST.get('recommendations', '')

        if photo_url and photo_url.startswith('/media/'):
            photo_path = photo_url.replace('/media/', '')
            full_path = os.path.join(settings.MEDIA_ROOT, photo_path)

            try:
                with open(full_path, 'rb') as f:
                    file_content = ContentFile(f.read(), name=os.path.basename(photo_path))

                SafetyObservation.objects.create(
                    photo=file_content,
                    hazard_description=hazard_description,
                    recommendations=recommendations,
                    created_by=request.user  # ✅ correct field
                )

                return redirect('mrsafe_app:inspect_success')  # ensure this route exists

            except FileNotFoundError:
                print(f"[ERROR] File not found: {full_path}")

        return redirect('mrsafe_app:inspect')  # fallback
    return redirect('mrsafe_app:inspect')


from django.shortcuts import render

def inspect_success(request):
    return render(request, "mrsafe/inspect/success.html")

from django.contrib.auth.decorators import login_required
from django.shortcuts import render
from ..models import SafetyObservation

@login_required
def observation_list(request):
    observations = SafetyObservation.objects.filter(created_by=request.user).order_by('-detected_at')
    return render(request, "mrsafe/inspect/observation_list.html", {
        "observations": observations
    })
###################################################################################

####################site inspection##########################################
# views/site_inspection_views.py
from django.shortcuts import render, get_object_or_404, redirect
from django.contrib.auth.decorators import login_required
from ..models import SiteInspection, SafetyObservation

@login_required
def inspection_list(request):
    inspections = SiteInspection.objects.filter(inspector=request.user).order_by('-date')
    return render(request, 'mrsafe/inspect/inspection_list.html', {
        'inspections': inspections
    })

from django.contrib.auth.decorators import login_required
from django.shortcuts import render, get_object_or_404
from ..models import SiteInspection, SafetyObservation

@login_required
def inspection_detail(request, inspection_id):
    # Ensure the user is the inspector of this SiteInspection
    inspection = get_object_or_404(SiteInspection, id=inspection_id, inspector=request.user)

    # Fetch related SafetyObservations and order by detection time
    observations = SafetyObservation.objects.filter(site_inspection=inspection).order_by('-detected_at')

    # Render the page with the new template and pass the inspection and observations context
    return render(request, 'mrsafe/inspect/inspection_detail.html', {
        'inspection': inspection,
        'observations': observations,
    })


from django.shortcuts import render, redirect
from django.contrib.auth.decorators import login_required
from ..forms import SiteInspectionForm  # make sure this form is defined
from ..models import SiteInspection
from django.urls import reverse


@login_required
def inspection_create(request):
    if request.method == 'POST':
        form = SiteInspectionForm(request.POST)
        if form.is_valid():
            inspection = form.save(commit=False)
            inspection.inspector = request.user
            inspection.save()
            return redirect('mrsafe_app:inspection_detail', inspection_id=inspection.id)


    else:
        form = SiteInspectionForm()
    return render(request, 'mrsafe/inspect/create.html', {'form': form})

from django.shortcuts import render, get_object_or_404
from ..models import SafetyObservation

def observation_detail(request, pk):
    # Fetch the SafetyObservation object
    obs = get_object_or_404(SafetyObservation, pk=pk)

    # Split the hazard description and recommendations
    hazard_lines = obs.hazard_description.split("\n") if obs.hazard_description else []
    reco_lines = obs.recommendations.split("\n") if obs.recommendations else []

    # Combine hazards and recommendations in a structured format
    hazards_and_recommendations = []
    for hazard, reco in zip(hazard_lines, reco_lines):
        hazards_and_recommendations.append({
            'hazard': hazard,
            'recommendation': reco
        })

    return render(request, "mrsafe/inspect/observation_detail.html", {
        "obs": obs,
        "hazards_and_recommendations": hazards_and_recommendations,
    })

import re
from django.views.decorators.csrf import csrf_exempt
from django.shortcuts import render, get_object_or_404
from django.conf import settings
from django.core.files.storage import FileSystemStorage
import base64
import traceback

@csrf_exempt
def site_inspection_image_test(request, inspection_id):
    context = {}
    inspection = get_object_or_404(SiteInspection, id=inspection_id)

    if not settings.OPENAI_API_KEY:
        context["error"] = "OpenAI API key is missing"
        return render(request, "mrsafe/inspect/site_inspect.html", context)

    if request.method == "POST" and request.FILES.get("photo"):
        image_file = request.FILES["photo"]

        try:
            # Save uploaded image
            fs = FileSystemStorage(location="media/uploads", base_url="/media/uploads/")
            filename = fs.save(image_file.name, image_file)
            uploaded_url = fs.url(filename)
            context["photo_url"] = uploaded_url

            # Encode for GPT
            with fs.open(filename, "rb") as img:
                base64_img = base64.b64encode(img.read()).decode("utf-8")
            content_type = image_file.content_type or "image/jpeg"
            image_ext = content_type.split("/")[-1]

            # GPT prompt
            prompt = """
**Safety Inspection Analysis Request**


You are an AI safety inspection assistant. Based on the provided input (such as a photo or observation), generate clearly separated and numbered hazard reports. 

Each hazard must be output using the following exact structure. Do not combine or merge hazards. Each block must be self-contained.

For every identified issue, follow **this markdown format**:

---

### Hazard #{{ n }}
- **Title**: [Short title of hazard]
- **Severity**: Low / Medium / High / Critical
- **Description**: [Detailed explanation of the hazard and its impact]
- **OSHA Reference**: [e.g., 1926.501(b)(1)] or "N/A" if not applicable

#### Recommendation
- **Action**: [Corrective or preventive measure to eliminate or reduce the hazard]

---

### ✅ Example Output

```markdown
### Hazard #1
- **Title**: Missing Guardrails on Elevated Platform
- **Severity**: High
- **Description**: The platform at the construction site lacks guardrails on two sides, increasing the risk of falls from height.
- **OSHA Reference**: 1926.501(b)(1)

#### Recommendation
- **Action**: Install compliant guardrails (at least 42 inches high) on all open sides of the platform immediately.

---

### Hazard #2
- **Title**: Blocked Emergency Exit
- **Severity**: Critical
- **Description**: The designated emergency exit is obstructed by equipment, making it inaccessible during an emergency.
- **OSHA Reference**: 1910.37(a)(3)

#### Recommendation
- **Action**: Clear the emergency exit path and place visible signage to ensure constant accessibility.


---

Repeat the structure for each hazard found in the image. Do not list hazards or recommendations in separate sections. Always keep each hazard and its recommendation grouped.

Format your output in clean markdown as shown above.
"""

            # GPT-4o call
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "You are a workplace safety AI assistant."},
                    {"role": "user", "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {
                            "url": f"data:image/{image_ext};base64,{base64_img}",
                            "detail": "high"
                        }}
                    ]}
                ],
                max_tokens=1500,
                temperature=0.3,
            )

            if not response.choices or not response.choices[0].message.content:
                raise ValueError("AI response is empty.")

            # Clean up response
            analysis = response.choices[0].message.content
            clean_analysis = analysis.replace("```markdown", "").replace("```", "").strip()

            # Extract hazard and recommendation sections
            def extract_sections(text):
                hazard_sections = []
                recommendation_sections = []
                blocks = text.split('---')
                for block in blocks:
                    block = block.strip()
                    if not block:
                        continue
                    hazard_match = re.search(r"###\s*Hazard.*?\n(.*?)(?=####|$)", block, re.DOTALL)
                    reco_match = re.search(r"####\s*Recommendation.*?\n(.*)", block, re.DOTALL)
                    if hazard_match:
                        hazard_sections.append(hazard_match.group(1).strip())
                    if reco_match:
                        recommendation_sections.append(reco_match.group(1).strip())
                return "\n\n".join(hazard_sections), "\n\n".join(recommendation_sections)

            hazard_text, recommendation_text = extract_sections(clean_analysis)

            # Pass to template
            context["hazard_description"] = hazard_text or clean_analysis
            context["recommendations"] = recommendation_text or "Review the analysis above and extract the key action items."
            context["inspection_id"] = inspection.id

            return render(request, "mrsafe/inspect/preview_observation.html", context)

        except Exception as e:
            context["error"] = f"Error during processing: {e}"
            if settings.DEBUG:
                context["traceback"] = traceback.format_exc()

    return render(request, "mrsafe/inspect/site_inspect.html", context)


# views.py
@csrf_exempt
def site_inspection_start(request, inspection_id):
    # Fetch the SiteInspection object using the provided inspection_id
    inspection = get_object_or_404(SiteInspection, id=inspection_id)

    if request.method == "POST" and request.FILES.get("photo"):
        image_file = request.FILES["photo"]

        try:
            # Validate the image size (max 5MB)
            if image_file.size > 5 * 1024 * 1024:
                raise ValueError("Image size exceeds 5MB")

            # Save the uploaded photo
            fs = FileSystemStorage(location="media/uploads", base_url="/media/uploads/")
            filename = fs.save(image_file.name, image_file)
            uploaded_url = fs.url(filename)

            # Now, we redirect to the image analysis view
            return redirect('mrsafe_app:site_inspection_image_test', inspection_id=inspection.id)

        except ValueError as ve:
            context["error"] = f"Validation Error: {str(ve)}"
        except Exception as e:
            context["error"] = f"Error during photo processing: {str(e)}"
            if settings.DEBUG:
                context["error_details"] = traceback.format_exc()

    return render(request, 'mrsafe/inspect/start_inspection.html', {'inspection': inspection})


from django.views.generic import TemplateView
from django.db.models import (
    Count, Avg, Q, F, ExpressionWrapper, 
    fields, Max, Case, When, Value, BooleanField
)
from django.views.generic import TemplateView
from django.db.models import (
    Count, Avg, Q, F, ExpressionWrapper, 
    fields, Case, When, Value, BooleanField
)
from django.db.models.functions import TruncMonth
from django.utils import timezone
from datetime import timedelta
from ..models import SiteInspection, SafetyObservation

class SafetyDashboardView(TemplateView):
    template_name = 'mrsafe/inspect/safety_dashboard.html'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        thirty_days_ago = timezone.now() - timedelta(days=30)
        current_user = self.request.user

        # 1. Inspection Metrics (only for current user)
        inspections = SiteInspection.objects.filter(
            inspector=current_user
        ).select_related('inspector')
        
        completed_inspections = inspections.filter(completed=True)
        
        # Ensure we always have data for the chart
        inspection_by_month = list(inspections.annotate(
            month=TruncMonth('date')
        ).values('month').annotate(
            count=Count('id')
        ).order_by('month'))
        
        context['inspection_metrics'] = {
            'total': inspections.count(),
            'completed': completed_inspections.count(),
            'completion_rate': round(
                (completed_inspections.count() / inspections.count() * 100) 
                if inspections.count() else 0, 
                1
            ),
            'avg_completion_time': completed_inspections.annotate(
                duration=ExpressionWrapper(
                    F('completed_at') - F('date'),
                    output_field=fields.DurationField()
                )
            ).aggregate(
                avg_duration=Avg('duration')
            )['avg_duration'],
            'by_month': inspection_by_month or [{
                'month': timezone.now().replace(day=1),
                'count': 0
            }]
        }

        # 2. Observation Metrics (only for current user)
        observations = SafetyObservation.objects.filter(
            created_by=current_user
        ).select_related('site_inspection')
        
        # Get hazard types with fallback
        hazard_types = list(observations.exclude(
            structured_json__hazard_type__isnull=True
        ).values(
            'structured_json__hazard_type'
        ).annotate(
            count=Count('id')
        ).order_by('-count')[:5])
        
        context['observation_metrics'] = {
            'total': observations.count(),
            'with_photos': observations.exclude(photo='').count(),
            'recent': observations.filter(detected_at__gte=thirty_days_ago).count(),
            'by_inspection': list(observations.values(
                'site_inspection__title'
            ).annotate(
                count=Count('id')
            ).order_by('-count')[:5]),
            'hazard_types': hazard_types or [{
                'structured_json__hazard_type': 'No hazards recorded', 
                'count': 0
            }]
        }

        # 3. Recent Activity (only for current user)
        context['recent_activity'] = {
            'inspections': list(inspections.order_by('-date')[:5]),
            'observations': list(observations.order_by('-detected_at')[:5])
        }

        # Debug output
        print("\n=== DASHBOARD DATA ===")
        print(f"User: {current_user.username}")
        print(f"Inspections: {context['inspection_metrics']['total']}")
        print(f"Observations: {context['observation_metrics']['total']}")
        print(f"Last Month: {context['observation_metrics']['recent']} observations")
        
        return context
    
from django.contrib.auth.decorators import login_required
from django.core.files.base import ContentFile
from django.shortcuts import redirect, get_object_or_404, render
from django.conf import settings
import os
from urllib.parse import unquote
from ..models import SafetyObservation, SiteInspection

@login_required
def safe_site_observation(request, inspection_id):
    if request.method == 'POST':
        photo_url = request.POST.get('photo_url')
        hazard_description = request.POST.get('hazard_description', '')
        recommendations = request.POST.get('recommendations', '')
        inspection = get_object_or_404(SiteInspection, id=inspection_id)

        if photo_url and photo_url.startswith('/media/'):
            photo_path = unquote(photo_url.replace('/media/', ''))
            full_path = os.path.join(settings.MEDIA_ROOT, photo_path)

            try:
                with open(full_path, 'rb') as f:
                    file_content = ContentFile(f.read(), name=os.path.basename(photo_path))

                SafetyObservation.objects.create(
                    photo=file_content,
                    hazard_description=hazard_description,
                    recommendations=recommendations,
                    created_by=request.user,
                    site_inspection=inspection
                )

                return redirect('mrsafe_app:inspection_detail', inspection_id=inspection.id)

            except FileNotFoundError:
                return render(request, "mrsafe/inspect/preview_observation.html", {
                    "error": f"File not found: {photo_url}",
                    "inspection_id": inspection.id,
                    "hazard_description": hazard_description,
                    "recommendations": recommendations,
                })

    return redirect('mrsafe_app:inspection_detail', inspection_id=inspection_id)



from django.utils.timezone import now
from django.contrib.auth.decorators import login_required
from django.shortcuts import get_object_or_404, redirect

@login_required
def finalize_inspection(request, inspection_id):
    inspection = get_object_or_404(SiteInspection, id=inspection_id, inspector=request.user)

    if request.method == "POST":
        inspection.completed = True
        inspection.completed_at = now()
        inspection.save()

    return redirect("mrsafe_app:inspection_list")


from django.template.loader import get_template
from django.templatetags.static import static
from django.http import HttpResponse
from xhtml2pdf import pisa
from django.contrib.auth.decorators import login_required
from ..models import SiteInspection
from django.shortcuts import get_object_or_404

@login_required
def export_inspection_pdf(request, inspection_id):
    inspection = get_object_or_404(SiteInspection, id=inspection_id)
    
    # Build full logo URL
    logo_url = request.build_absolute_uri(static('images/logo-company4.png'))

    # Render HTML with logo_url
    template = get_template('mrsafe/inspect/inspection_pdf.html')
    html = template.render({'inspection': inspection, 'logo_url': logo_url})

    # Create PDF
    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = f'inline; filename="inspection_{inspection.id}.pdf"'

    pisa_status = pisa.CreatePDF(html, dest=response)

    if pisa_status.err:
        return HttpResponse("Error generating PDF", status=500)
    
    return response  # ✅ Return PDF, not render()


from django.contrib.auth.decorators import login_required
from django.shortcuts import get_object_or_404
from django.http import HttpResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
from datetime import datetime
from ..models import SiteInspection

@login_required
def export_inspection_ppt(request, inspection_id):
    inspection = get_object_or_404(SiteInspection, id=inspection_id)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    PRIMARY_COLOR = RGBColor(44, 62, 80)
    SECONDARY_COLOR = RGBColor(52, 152, 219)
    ACCENT_COLOR = RGBColor(231, 76, 60)
    LIGHT_BG = RGBColor(245, 245, 245)

    # ---------------- Title Slide ----------------
    blank_layout = prs.slide_layouts[6]
    title_slide = prs.slides.add_slide(blank_layout)

    # Background
    background = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR

    # Logo (use actual file path or catch error)
    logo_path = 'static/images/mrsafe-logo.png'  # Adjust path as needed
    try:
        title_slide.shapes.add_picture(
            logo_path, Inches(1), Inches(0.5), height=Inches(1)
        )
    except:
        logo_box = title_slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(2), Inches(0.5))
        logo_text = logo_box.text_frame.add_paragraph()
        logo_text.text = "MrSafe"
        logo_text.font.size = Pt(24)
        logo_text.font.color.rgb = RGBColor(255, 255, 255)

    # Title
    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(2), prs.slide_width - Inches(2), Inches(1.5))
    title = title_box.text_frame.add_paragraph()
    title.text = "Site Inspection Report"
    title.font.size = Pt(44)
    title.font.color.rgb = RGBColor(255, 255, 255)
    title.font.bold = True

    # Info
    info_box = title_slide.shapes.add_textbox(Inches(1), Inches(3.5), prs.slide_width - Inches(2), Inches(3))
    info = info_box.text_frame
    info.word_wrap = True

    def safe_name(val):
        return val.get_full_name() if hasattr(val, 'get_full_name') else str(val)

    info_lines = [
        f"Inspection ID: MR-{inspection.id:06d}",
        f"Title: {inspection.title}",
        f"Location: {inspection.location}",
        f"Inspector: {safe_name(getattr(inspection, 'inspector', 'N/A'))}",
    ]

    for line in info_lines:
        p = info.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Footer
    footer_box = title_slide.shapes.add_textbox(Inches(1), prs.slide_height - Inches(0.7), prs.slide_width - Inches(2), Inches(0.5))
    footer = footer_box.text_frame.add_paragraph()
    footer.text = "Generated by MrSafe | www.mrsafe.me"
    footer.font.size = Pt(12)
    footer.font.color.rgb = RGBColor(200, 200, 200)
    footer.alignment = PP_ALIGN.RIGHT

    # ---------------- Summary Slide ----------------
    summary_slide = prs.slides.add_slide(prs.slide_layouts[5])
    summary_slide.shapes.title.text = "Inspection Summary"
    summary_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    summary_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)

    # Background
    bg = summary_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    summary_slide.shapes._spTree.remove(bg._element)
    summary_slide.shapes._spTree.insert(2, bg._element)

    rows, cols = 4, 2
    table = summary_slide.shapes.add_table(
        rows, cols, Inches(1.5), Inches(2), prs.slide_width - Inches(3), Inches(3.5)
    ).table

    table.columns[0].width = Inches(3)
    table.columns[1].width = prs.slide_width - Inches(3) - Inches(3)

    # Headers
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    for col in range(2):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER

    # Data rows
    data = [
        ("Total Observations", str(inspection.observations.count())),
        ("High Risk Items", str(getattr(inspection, 'high_risk_count', 0))),
        ("Completion", f"{getattr(inspection, 'completion_percentage', 0)}%"),
    ]

    for row_idx, (label, value) in enumerate(data, start=1):
        table.cell(row_idx, 0).text = label
        table.cell(row_idx, 1).text = value

        for col in range(2):
            cell = table.cell(row_idx, col)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # ---------------- Observations Slides ----------------
    observations = inspection.observations.all()
    if observations.exists():
        for idx, obs in enumerate(observations, 1):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Observation #{idx}"
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)

            # Background
            obs_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            obs_bg.fill.solid()
            obs_bg.fill.fore_color.rgb = LIGHT_BG
            slide.shapes._spTree.remove(obs_bg._element)
            slide.shapes._spTree.insert(2, obs_bg._element)

            # Textbox
            box = slide.shapes.add_textbox(Inches(1), Inches(1.8), prs.slide_width - Inches(2), prs.slide_height - Inches(2.5))
            tf = box.text_frame
            tf.word_wrap = True

            # Hazard
            p1 = tf.add_paragraph()
            p1.text = "HAZARD:"
            p1.font.size = Pt(16)
            p1.font.color.rgb = ACCENT_COLOR
            p1.font.bold = True

            p2 = tf.add_paragraph()
            p2.text = obs.hazard_description
            p2.font.size = Pt(14)

            # Recommendation
            p3 = tf.add_paragraph()
            p3.text = "RECOMMENDATION:"
            p3.font.size = Pt(16)
            p3.font.color.rgb = SECONDARY_COLOR
            p3.font.bold = True

            p4 = tf.add_paragraph()
            p4.text = obs.recommendations
            p4.font.size = Pt(14)

            # Footer with detection date
            footer_box = slide.shapes.add_textbox(Inches(1), prs.slide_height - Inches(0.7), prs.slide_width - Inches(2), Inches(0.5))
            footer = footer_box.text_frame.add_paragraph()
            footer.text = f"Detected on: {obs.detected_at.strftime('%B %d, %Y') if getattr(obs, 'detected_at', None) else 'Date not specified'}"
            footer.font.size = Pt(12)
            footer.font.color.rgb = RGBColor(150, 150, 150)
    else:
        # No Observations Slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "No Observations Recorded"
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)

        box = slide.shapes.add_textbox(Inches(2), Inches(3), prs.slide_width - Inches(4), Inches(2))
        p = box.text_frame.add_paragraph()
        p.text = "No safety observations were recorded for this inspection."
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    # ---------------- Export as Response ----------------
    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)

    date_str = datetime.today().strftime("%Y%m%d")
    filename = f"mrsafe_inspection_{inspection.id}_{date_str}.pptx"

    response = HttpResponse(
        ppt_io.read(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    response['Content-Disposition'] = f'attachment; filename="{filename}"'

    return response



from django.contrib.auth.decorators import login_required
from django.shortcuts import render, get_object_or_404
from ..models import SiteInspection

@login_required
def inspection_full_report(request, inspection_id):
    inspection = get_object_or_404(SiteInspection, id=inspection_id)

    return render(request, "mrsafe/inspect/full_report.html", {
        "inspection": inspection,
        "observations": inspection.observations.all(),
    })

from django.contrib.auth.decorators import login_required
from django.shortcuts import get_object_or_404
from django.http import HttpResponse
from ..models import SiteInspection

@login_required
def export_inspection_docx(request, inspection_id):
    inspection = get_object_or_404(SiteInspection, id=inspection_id)

    # Temporary placeholder until export logic is implemented
    return HttpResponse(f"DOCX export for inspection: {inspection.title}", content_type="text/plain")
